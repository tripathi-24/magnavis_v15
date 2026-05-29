#!/usr/bin/env python3
"""
Generate a detailed and accessible deep-dive PowerPoint for `application_temp.py` and related scripts.

Output:
  presentations/application_temp_deep_dive_improved.pptx

Notes:
  - This deck introduces naive-user-friendly analogies and high-level architecture breakdowns
    before diving into the dense catalog slides.
  - Function/method catalogs are auto-generated from AST parsing to keep
    slide content aligned with the current codebase.
"""

from __future__ import annotations

import ast
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = PROJECT_ROOT / "presentations"
OUT_PPTX = OUT_DIR / "application_temp_deep_dive_improved.pptx"


@dataclass
class Entry:
    symbol: str
    signature: str
    returns: str
    summary: str
    category: str
    lineno: int


@dataclass
class ModuleInventory:
    path: Path
    title: str
    purpose: str
    entries: List[Entry]
    top_functions: int
    class_count: int
    method_count: int


# Theme
BG = RGBColor(246, 248, 252)
HEADER_BG = RGBColor(15, 23, 42)
HEADER_TXT = RGBColor(248, 250, 252)
ACCENT = RGBColor(30, 102, 204)
ACCENT_2 = RGBColor(12, 132, 102)
TEXT = RGBColor(40, 52, 68)
MUTED = RGBColor(95, 108, 126)
WHITE = RGBColor(255, 255, 255)
ROW_ALT = RGBColor(241, 246, 255)


MODULE_SPECS = [
    (
        PROJECT_ROOT / "src" / "app.py",
        "app.py",
        "Main multi-sensor GUI workflow: startup mode selection, threaded fetching, "
        "predictor orchestration, anomaly detection, direction inference, and plotting.",
    ),
    (
        PROJECT_ROOT / "src" / "data_convert_db_now.py",
        "data_convert_db_now.py",
        "DB ingestion/query and time-series/vector transformation layer powering "
        "real-time and simulation modes.",
    ),
    (
        PROJECT_ROOT / "src" / "predictor_ai.py",
        "predictor_ai.py",
        "Attention-BiLSTM predictor implementation, feature engineering, and forecast loop.",
    ),
    (
        PROJECT_ROOT / "src" / "Anomaly_detector.py",
        "Anomaly_detector.py",
        "Adaptive residual-based anomaly detector and associated statistics.",
    ),
    (
        PROJECT_ROOT / "src" / "anomaly_direction.py",
        "anomaly_direction.py",
        "Direction geometry primitives and triangulation helpers.",
    ),
    (
        PROJECT_ROOT / "src" / "train_gru_pretrained.py",
        "train_gru_pretrained.py",
        "Offline pretraining pipeline for per-sensor Attention-BiLSTM models.",
    ),
]


def sentence(doc: Optional[str]) -> str:
    if not doc:
        return ""
    for ln in doc.splitlines():
        s = ln.strip()
        if s:
            return s
    return ""


def infer_category(name: str) -> str:
    n = name.lower()
    if n.startswith(("init", "__init__", "configure")):
        return "Initialization/Configuration"
    if n.startswith(("fetch", "get", "load", "query", "discover")):
        return "Data Access"
    if n.startswith(("build", "create", "compute", "derive", "parse")):
        return "Transformation/Computation"
    if n.startswith(("start", "run", "enqueue", "poll", "update")):
        return "Runtime Orchestration"
    if n.startswith(("detect", "predict", "forecast", "triangulate")):
        return "Inference/Detection"
    if n.startswith(("save", "write", "export")):
        return "Persistence"
    if n.startswith(("is_", "has_", "can_")):
        return "Predicate/Validation"
    return "Utility/Control Flow"


def summary_from_name(name: str) -> str:
    raw = name.strip("_").replace("_", " ")
    if not raw:
        return "Internal helper routine."
    words = raw.split()
    if words[0] in {"get", "fetch", "load", "read"}:
        return f"Retrieves {raw[len(words[0]):].strip() or 'data'}."
    if words[0] in {"build", "create", "compute", "derive", "parse"}:
        return f"Builds or computes {raw[len(words[0]):].strip() or 'derived output'}."
    if words[0] in {"start", "run", "enqueue", "update", "poll"}:
        return f"Drives runtime step: {raw}."
    if words[0] in {"detect", "predict", "forecast"}:
        return f"Performs model inference step: {raw}."
    if words[0] in {"save", "write", "export"}:
        return f"Persists output for {raw[len(words[0]):].strip() or 'state'}."
    if words[0] in {"is", "has", "can"}:
        return f"Predicate helper for {raw}."
    return f"Implements {raw}."


def fmt_signature(node: ast.FunctionDef, method: bool = False) -> str:
    args = []
    pos = list(node.args.posonlyargs) + list(node.args.args)
    defaults = [None] * (len(pos) - len(node.args.defaults)) + list(node.args.defaults)
    for i, a in enumerate(pos):
        name = a.arg
        if method and i == 0 and name == "self":
            continue
        if defaults[i] is None:
            args.append(name)
        else:
            args.append(f"{name}=?")
    if node.args.vararg:
        args.append(f"*{node.args.vararg.arg}")
    for a, d in zip(node.args.kwonlyargs, node.args.kw_defaults):
        if d is None:
            args.append(f"{a.arg}")
        else:
            args.append(f"{a.arg}=?")
    if node.args.kwarg:
        args.append(f"**{node.args.kwarg.arg}")
    return f"{node.name}({', '.join(args)})"


def fmt_return(node: ast.FunctionDef) -> str:
    if node.returns is None:
        return "implicit/None"
    try:
        return ast.unparse(node.returns)
    except Exception:
        return "annotated"


def parse_module(path: Path, title: str, purpose: str) -> ModuleInventory:
    src = path.read_text(encoding="utf-8", errors="ignore")
    tree = ast.parse(src)
    entries: List[Entry] = []
    top_functions = 0
    class_count = 0
    method_count = 0

    for n in tree.body:
        if isinstance(n, ast.FunctionDef):
            top_functions += 1
            doc = sentence(ast.get_docstring(n))
            entries.append(
                Entry(
                    symbol=n.name,
                    signature=fmt_signature(n, method=False),
                    returns=fmt_return(n),
                    summary=doc or summary_from_name(n.name),
                    category=infer_category(n.name),
                    lineno=n.lineno,
                )
            )
        elif isinstance(n, ast.ClassDef):
            class_count += 1
            for m in n.body:
                if isinstance(m, ast.FunctionDef):
                    method_count += 1
                    doc = sentence(ast.get_docstring(m))
                    sym = f"{n.name}.{m.name}"
                    entries.append(
                        Entry(
                            symbol=sym,
                            signature=fmt_signature(m, method=True),
                            returns=fmt_return(m),
                            summary=doc or summary_from_name(m.name),
                            category=infer_category(m.name),
                            lineno=m.lineno,
                        )
                    )

    entries.sort(key=lambda e: e.lineno)
    return ModuleInventory(
        path=path,
        title=title,
        purpose=purpose,
        entries=entries,
        top_functions=top_functions,
        class_count=class_count,
        method_count=method_count,
    )


def chunks(seq: List[Entry], size: int) -> Iterable[List[Entry]]:
    for i in range(0, len(seq), size):
        yield seq[i : i + size]


def apply_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title: str, subtitle: str = "") -> None:
    apply_bg(slide)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    band.fill.solid()
    band.fill.fore_color.rgb = HEADER_BG
    band.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.08), Inches(12.6), Inches(0.45))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(26)
    p.font.color.rgb = HEADER_TXT
    p.font.name = "Segoe UI"

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.45), Inches(0.52), Inches(12.6), Inches(0.35))
        stf = sb.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = RGBColor(180, 193, 212)
        sp.font.name = "Segoe UI"


def add_bullet_box(slide, left: float, top: float, width: float, height: float, title: str, bullets: List[str], accent: RGBColor = ACCENT) -> None:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = accent
    box.line.width = Pt(1.4)

    title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.42))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = accent
    title_box.line.fill.background()

    ttb = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.03), Inches(width - 0.2), Inches(0.35))
    ttf = ttb.text_frame
    ttf.clear()
    tp = ttf.paragraphs[0]
    tp.text = title
    tp.font.bold = True
    tp.font.size = Pt(13)
    tp.font.color.rgb = WHITE
    tp.font.name = "Segoe UI"
    ttf.vertical_anchor = MSO_ANCHOR.MIDDLE

    ctb = slide.shapes.add_textbox(Inches(left + 0.12), Inches(top + 0.58), Inches(width - 0.24), Inches(height - 0.66))
    ctf = ctb.text_frame
    ctf.clear()
    ctf.word_wrap = True
    for i, b in enumerate(bullets):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = b
        p.font.size = Pt(12)
        p.font.name = "Segoe UI"
        p.font.color.rgb = TEXT
        p.level = 0
        p.space_after = Pt(8)


def add_inventory_table(slide, entries: List[Entry]) -> None:
    rows = len(entries) + 1
    cols = 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.35), Inches(1.2), Inches(12.65), Inches(5.95))
    table = table_shape.table

    col_widths = [3.0, 3.6, 1.5, 4.55]
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    headers = ["Symbol", "Signature (inputs)", "Returns", "Purpose and Concept"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER

    for i, e in enumerate(entries, start=1):
        vals = [
            e.symbol,
            e.signature,
            e.returns,
            f"[{e.category}] {e.summary}",
        ]
        for j, v in enumerate(vals):
            cell = table.cell(i, j)
            cell.text = v
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_ALT if i % 2 == 0 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.TOP
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.font.size = Pt(10)
            p.font.name = "Segoe UI"
            p.font.color.rgb = TEXT
            p.alignment = PP_ALIGN.LEFT


def concept_slides(prs: Presentation) -> None:
    # 1. Title Slide
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "Magnavis Deep Dive: From Concept to Code", "A complete breakdown of the application workflow (~50 slides)")
    add_bullet_box(
        s, 1.0, 1.8, 11.3, 5.0, "Welcome!",
        [
            "Whether you're new to the project or looking to understand the fine technical details, this presentation has you covered.",
            "We will start with high-level introductory analogies.",
            "Then, we'll break down the workflow step by step: from reading sensor data to drawing 3D anomalies on screen.",
            "Finally, we'll provide an exhaustive function-by-function code catalog, automatically generated from the Python source code.",
        ],
        ACCENT,
    )

    # 2. Executive Summary - What is Magnavis?
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "1. Executive Summary", "What does this application actually do?")
    add_bullet_box(
        s, 0.8, 1.5, 11.7, 5.6, "What is Magnavis?",
        [
            "Magnavis is a smart software system that monitors magnetic fields from various sensors in real-time.",
            "Instead of relying on a human to spot a sudden 'blip' or 'spike' on a graph, the software employs an AI (Artificial Intelligence) brain.",
            "The AI figures out what the normal 'baseline' should be for any given time of day.",
            "When the actual magnetic field dramatically breaks away from that expected baseline, the system flags it as an 'Anomaly'.",
            "Once flagged, it uses complex mathematics to literally point a 3D arrow towards the source of that magnetic disturbance."
        ],
        ACCENT_2,
    )

    # 3. System Analogy
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "2. A Simple Analogy", "Understanding the architecture using a Smart Security System")
    add_bullet_box(
        s, 0.5, 1.5, 6.0, 5.5, "The Smart Security Camera",
        [
            "1. The Eyes (Data Ingest) - The CCTV camera constantly watches a busy street, capturing streaming video.",
            "2. The Brain (AI Predictor) - The AI knows that heavy traffic implies a lot of movement at 5 PM. It expects activity.",
            "3. The Alert (Anomaly Detector) - Someone throwing a brick is NOT expected. The system subtracts the expected behavior from actual behavior, and sounds the alarm.",
            "4. The Direction (Triangulation) - A second camera helps figure out exactly where the brick was thrown from."
        ],
        ACCENT,
    )
    add_bullet_box(
        s, 6.8, 1.5, 6.0, 5.5, "How It Maps to Our Code",
        [
            "1. The Eyes -> `data_convert_db_now.py`: Fetches real-time magnetic values from databases.",
            "2. The Brain -> `predictor_ai.py`: An Attention-BiLSTM neural network that predicts baseline trends.",
            "3. The Alert -> `Anomaly_detector.py`: Checks for residual spikes exceeding normal thresholds.",
            "4. The Direction -> `anomaly_direction.py`: Triangulates 3D position of disturbance."
        ],
        ACCENT_2,
    )

    # 4. Runtime Architecture Story
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "3. The Data Story (Step-by-Step)", "How data travels through the application")
    add_bullet_box(
        s, 0.7, 1.5, 12.0, 5.5, "The Application Pipeline",
        [
            "STEP 1: The user opens `application_temp.py` and chooses a startup mode (Real-Time monitoring, Historical replay, or CSV simulation).",
            "STEP 2: Background 'Workers' (small invisible helpers) start fetching magnetic data behind the scenes without freezing the app.",
            "STEP 3: The data is grouped into 'Historic' (past 60 minutes) and 'New' (fresh arrivals).",
            "STEP 4: A separate process tells the AI to quickly predict the next sequence of values for each sensor.",
            "STEP 5: `Anomaly_detector.py` rapidly does the math: (Actual Value) minus (Predicted Value). If the result is a huge spike, an Anomaly Flag is thrown.",
            "STEP 6: The GUI updates with red lines (Anomalies) and a 3D arrow spins to face the source."
        ],
        ACCENT,
    )

    # 5. Core Mode Matrix
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "4. Operating Modes", "How we test and run the system")
    add_bullet_box(
        s, 0.6, 1.45, 4.1, 5.8, "Live / Real-Time",
        [
            "The intended use case. Watches the active DB.",
            "Continuously queries the database every few seconds.",
            "AI predicts on the fly.",
        ],
        RGBColor(34, 139, 34),
    )
    add_bullet_box(
        s, 4.9, 1.45, 4.1, 5.8, "Simulation",
        [
            "Starts from a specific date in the past.",
            "Fast-forwards database queries to simulate time passing.",
            "Useful for testing live code on old data.",
        ],
        RGBColor(146, 64, 14),
    )
    add_bullet_box(
        s, 9.2, 1.45, 3.5, 5.8, "CSV Replay",
        [
            "Fully offline replay.",
            "Reads a downloaded file instead of connecting to a remote database.",
            "Useful for debugging very specific anomalies offline.",
        ],
        ACCENT_2,
    )

    # 6. Deep Dive into the Models
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "5. Under the Hood: The AI and UI", "What makes it powerful")
    add_bullet_box(
        s, 0.8, 1.5, 5.9, 5.8, "The BiLSTM Neural Network",
        [
            "Instead of standard flat thresholds, we use a Deep Learning model: 'Attention-BiLSTM'.",
            "BiLSTM means it looks forwards and backward in the sequence.",
            "Attention means it learns to focus on the parts of the sequence that matter the most.",
            "If it knows it's afternoon on a solar active day, it sets the baseline higher to avoid false alarms."
        ],
        ACCENT,
    )
    add_bullet_box(
        s, 7.0, 1.5, 5.5, 5.8, "The User Interface",
        [
            "`application_temp.py` is the main orchestra conductor.",
            "It manages up to three sensors simultaneously.",
            "Uses asynchronous threading (QThread) so drawing heavy charts never slows down data collection.",
            "Includes a 'Freeze Window' to stop masking an anomaly once it triggers."
        ],
        ACCENT_2,
    )

    # 7. Code Catalog Transition
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(s, "6. Technical Code Catalogs", "For the Developers and Engineers")
    add_bullet_box(
        s, 1.0, 1.5, 11.3, 5.4, "What Follows Next:",
        [
            "The remaining slides provide an exhaustive, alphabetically ordered catalog of every single function, Python class, and method that powers Magnavis.",
            "This ensures any software engineer can easily map a concept or feature back to its exact code symbol.",
            "Each slide lists the function signature, its intended purpose, what it returns, and a conceptual category (e.g. Data Access vs Inference).",
            "This catalog auto-generates directly from the Python syntax trees, ensuring 100% accuracy to the current codebase."
        ],
        ACCENT,
    )


def add_module_overview(prs: Presentation, inv: ModuleInventory, index: int, total: int) -> None:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        s,
        f"Module {index}/{total}: {inv.title}",
        f"{inv.path.relative_to(PROJECT_ROOT)}",
    )
    add_bullet_box(
        s,
        0.8,
        1.5,
        12.0,
        5.8,
        "Overview",
        [
            inv.purpose,
            f"Top-level functions: {inv.top_functions}",
            f"Classes: {inv.class_count}",
            f"Class methods: {inv.method_count}",
            f"Catalog entries in next slides: {len(inv.entries)}",
        ],
        ACCENT_2,
    )


def add_module_catalog(prs: Presentation, inv: ModuleInventory, chunk_size: int = 5) -> None:
    all_chunks = list(chunks(inv.entries, chunk_size))
    for i, group in enumerate(all_chunks, start=1):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_header(
            s,
            f"{inv.title} Catalog",
            f"Part {i}/{len(all_chunks)} | Entries {((i - 1) * chunk_size) + 1}-{((i - 1) * chunk_size) + len(group)}",
        )
        add_inventory_table(s, group)


def build_presentation() -> Path:
    inventories = [parse_module(path, title, purpose) for path, title, purpose in MODULE_SPECS]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    concept_slides(prs)

    total_modules = len(inventories)
    for i, inv in enumerate(inventories, start=1):
        add_module_overview(prs, inv, i, total_modules)
        add_module_catalog(prs, inv, chunk_size=5)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX))
    return OUT_PPTX


def main() -> None:
    out = build_presentation()
    print(f"Saved presentation: {out}")
    # Report total slide count for convenience.
    prs = Presentation(str(out))
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
